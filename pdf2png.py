from pdf2image import convert_from_path
from PIL import Image, ImageDraw, ImageFilter
import os
import glob
import img2pdf
from pptx import Presentation


# pdf to png
pdf_data = glob.glob('pdf/*.pdf')

img_path = 'img'
out_path = 'output'

if not os.path.exists(img_path):
    os.mkdir(img_path)

if not os.path.exists(out_path):
    os.mkdir(out_path)

for i in range(len(pdf_data)):
    basename = os.path.basename(pdf_data[i])
    file_name = os.path.splitext(basename)[0]

    convert_from_path(pdf_data[i], output_folder=img_path, fmt='png', output_file=file_name, single_file=False)


# png to pptx
prs = Presentation()
png_data = glob.glob('img/*.png')
slide_path = '{}/output.pptx'.format(out_path)

blank_slide_layout = prs.slide_layouts[6]

for i in range(len(png_data)):
    basename = os.path.basename(png_data[i])
    #file_name = os.path.splitext(basename)[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(basename, 0, 0)

prs.save(slide_path)